import subprocess
import sys

def install_package(package): #Instala el paquete necesario para leer y escribir en PowerPoint
    try:
        __import__(package)
    except ImportError:
        print(f"Instalando el paquete {package}...")
        subprocess.check_call([sys.executable, "-m", "pip", "install", package])

install_package("python-pptx")

#inicia el script
import os
from pptx import Presentation
from pptx.util import Inches

def create_pptx(file_name="trial.pptx"): #Crea un archivo PowerPoint con una diapositiva y un texto inicial "Hola". (En caso de que no esté previamente creado)    
    try:
        presentation = Presentation()
        slide = presentation.slides.add_slide(presentation.slide_layouts[5])
        textbox = slide.shapes.add_textbox(Inches(2), Inches(2), Inches(4), Inches(1))
        textbox.text = "Hola"
        presentation.save(file_name)
        print(f"Archivo '{file_name}' creado con éxito.")
    except Exception as e:
        print(f"Error al crear el archivo PowerPoint: {e}")

def replace_text_in_pptx(input_file, output_file, old_text, new_text): 
    try:
        if not os.path.exists(input_file): #Si el archivo por alguna razón no existe, lo creamos desde este script
            print(f"El archivo '{input_file}' no existe.")
            print("1. Crear archivo")
            print("2. Cancelar")
            create_option = input("Seleccione una opción: ").strip()
            if create_option == '1':
                create_pptx(input_file)
            else:
                print("Operación cancelada por el usuario.")
                return

        presentation = Presentation(input_file)
        replaced = False

        for slide in presentation.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    if old_text in shape.text:
                        shape.text = shape.text.replace(old_text, new_text)
                        replaced = True

        if replaced: #Se dan opciones
            print("1. Reemplazar archivo existente") #Con la primera (1) se sobrescribe el documento original
            print("2. Crear un nuevo archivo") #Con la segunda (2) se crea un nuevo documento con la modificación hecha
            save_option = input("Seleccione una opción: ").strip()
            if save_option == '1':
                presentation.save(input_file)
                print(f"El archivo '{input_file}' fue actualizado con éxito.")
                validate_replacement(input_file, new_text)
            elif save_option == '2':
                presentation.save(output_file)
                print(f"El texto se reemplazó correctamente. Archivo guardado como '{output_file}'.")
                validate_replacement(output_file, new_text)
            else:
                print("Opción no válida. No se realizaron cambios.")
        else:
            print(f"No se encontró el texto '{old_text}' en el archivo '{input_file}'.")
    except Exception as e:
        print(f"Error al reemplazar el texto: {e}")

def validate_replacement(file_path, expected_text): #Finalmente valida si el texto esperado está presente en el archivo PowerPoint, "Gracias" en este caso.
    try:
        if not os.path.exists(file_path):
            raise FileNotFoundError(f"El archivo '{file_path}' no existe.")

        presentation = Presentation(file_path)
        for slide in presentation.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    if expected_text in shape.text:
                        print(f"Validación exitosa: se encontró el texto '{expected_text}' en el archivo.")
                        return
        print(f"Validación fallida: no se encontró el texto '{expected_text}'.")
    except Exception as e:
        print(f"Error al validar el archivo: {e}")

if __name__ == "__main__":
    try:
        replace_text_in_pptx("trial.pptx", "trial_actualizado.pptx", "Hola", "Gracias")

    except Exception as e:
        print(f"Error al ejecutar el script: {e}")
